#!/usr/bin/env python3
"""
Certichain Africa — Investor Pitch Deck Generator
Output: Desktop/Certichain_PitchDeck.pptx  +  PDF (via PowerPoint COM if available)
"""

import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ─── Palette ──────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x1E, 0x52)
TEAL   = RGBColor(0x2D, 0xD4, 0xBF)
GOLD   = RGBColor(0xC9, 0xA8, 0x4C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF1, 0xF5, 0xF9)
MGRAY  = RGBColor(0x94, 0xA3, 0xB8)
DGRAY  = RGBColor(0x1E, 0x29, 0x3B)
BGRAY  = RGBColor(0xF8, 0xFA, 0xFC)

# ─── Canvas ───────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ─── Helpers ──────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)

def box(sl, x, y, w, h, fill):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    f = s.fill; f.solid(); f.fore_color.rgb = fill
    return s

def oval(sl, x, y, w, h, fill):
    s = sl.shapes.add_shape(9, x, y, w, h)
    s.line.fill.background()
    f = s.fill; f.solid(); f.fore_color.rgb = fill
    return s

def t(sl, text, x, y, w, h,
      size=14, bold=False, color=WHITE,
      align=PP_ALIGN.LEFT, italic=False, font='Calibri'):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    f   = run.font
    f.name = font; f.size = Pt(size); f.bold = bold
    f.italic = italic; f.color.rgb = color
    return tb

def tlines(sl, lines, x, y, w, h,
           size=13, bold=False, color=WHITE,
           align=PP_ALIGN.LEFT, font='Calibri', spc=6):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(spc)
        if isinstance(ln, dict):
            run = p.add_run(); run.text = ln.get('text','')
            f = run.font; f.name = font
            f.size  = Pt(ln.get('size', size))
            f.bold  = ln.get('bold', bold)
            f.italic= ln.get('italic', False)
            f.color.rgb = ln.get('color', color)
        else:
            run = p.add_run(); run.text = str(ln)
            f = run.font; f.name = font
            f.size = Pt(size); f.bold = bold; f.color.rgb = color
    return tb

def rule(sl, x, y, w, color=TEAL, h=Pt(4)):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    f = s.fill; f.solid(); f.fore_color.rgb = color
    return s

def header(sl, eyebrow, title, title_size=24):
    box(sl, 0, 0, W, Inches(1.5), NAVY)
    t(sl, eyebrow, Inches(0.6), Inches(0.18), Inches(10), Inches(0.4),
      size=10, bold=True, color=TEAL)
    t(sl, title,   Inches(0.6), Inches(0.55), Inches(12.1), Inches(0.9),
      size=title_size, bold=True, color=WHITE)

def footnote(sl, text):
    t(sl, text, Inches(0.55), Inches(7.1), Inches(12.2), Inches(0.38),
      size=8.5, color=MGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, NAVY)
box(sl, 0, 0, Inches(0.38), H, TEAL)               # left accent bar
oval(sl, W-Inches(4.2), -Inches(1.8), Inches(5.5), Inches(5.5),
     RGBColor(0x1A, 0x3A, 0x72))                   # top-right circle glow
oval(sl, -Inches(1.8), H-Inches(2.5), Inches(3.5), Inches(3.5),
     RGBColor(0x14, 0x30, 0x68))                   # bottom-left circle glow

# Logo tile
box(sl, Inches(0.9), Inches(1.8), Inches(0.72), Inches(0.72), TEAL)
t(sl, "C", Inches(0.9), Inches(1.79), Inches(0.72), Inches(0.72),
  size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Company name
t(sl, "Certichain Africa",
  Inches(0.87), Inches(2.65), Inches(7), Inches(1.0),
  size=46, bold=True, color=WHITE)

rule(sl, Inches(0.87), Inches(3.8), Inches(4.2), TEAL, Pt(4))

t(sl, "La certification numérique de confiance pour l'Afrique",
  Inches(0.87), Inches(3.97), Inches(8.5), Inches(0.9),
  size=19, color=RGBColor(0xCB, 0xD5, 0xE1))

t(sl, "Blockchain Polygon  ·  IPFS  ·  SHA-256  ·  Vérification instantanée",
  Inches(0.87), Inches(5.1), Inches(9), Inches(0.45),
  size=12, color=TEAL)

t(sl, "[ÉVÈNEMENT / DATE]",
  Inches(0.87), Inches(6.45), Inches(6), Inches(0.42),
  size=11, color=MGRAY, italic=True)

t(sl, "Confidentiel — Usage interne investisseurs",
  Inches(0.87), Inches(6.9), Inches(8), Inches(0.38),
  size=9, color=RGBColor(0x47, 0x55, 0x69))


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, WHITE)
header(sl, "LE PROBLÈME",
       "La fraude aux diplômes coûte des millions à l'Afrique chaque année")

stats = [
    ("30 %",    "des diplômes sont\nfrauduleux",
     "dans certains pays africains",
     RGBColor(0xFE, 0xE2, 0xE2), RGBColor(0x7F, 0x1D, 0x1D)),
    ("3–7 jours","pour vérifier\nun diplôme",
     "délai moyen employeur\n→ université",
     RGBColor(0xFF, 0xF7, 0xED), RGBColor(0x78, 0x35, 0x0F)),
    ("Coûts\nélevés", "impression · postal\narchivage physique",
     "gestion documentaire\ntraditionnelle",
     RGBColor(0xEF, 0xF6, 0xFF), RGBColor(0x1E, 0x1B, 0x4B)),
]
for i, (stat, line1, line2, bg, accent) in enumerate(stats):
    bx = Inches(0.55) + i * Inches(4.22)
    bw = Inches(4.0)
    box(sl, bx, Inches(1.78), bw, Inches(4.5), bg)
    box(sl, bx, Inches(1.78), bw, Pt(5), accent)
    t(sl, stat, bx+Inches(0.3), Inches(2.0), bw-Inches(0.5), Inches(1.0),
      size=38, bold=True, color=accent)
    t(sl, line1, bx+Inches(0.3), Inches(3.1), bw-Inches(0.5), Inches(0.75),
      size=14, bold=True, color=DGRAY)
    t(sl, line2, bx+Inches(0.3), Inches(3.85), bw-Inches(0.5), Inches(0.8),
      size=11, color=MGRAY)

footnote(sl, "* Sources estimatives — UNESCO Institute for Statistics, Banque Africaine de Développement [PLACEHOLDER]")


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, WHITE)
header(sl, "LA SOLUTION",
       "Certichain : des diplômes infalsifiables, vérifiables en 3 secondes")

points = [
    ("🔐", "Inviolable",
     "Chaque certificat est ancré sur Polygon. Toute modification est détectable instantanément."),
    ("⚡", "Instantané",
     "Vérification via QR code en 3 secondes. Aucun intermédiaire. Aucun formulaire."),
    ("🌍", "Africain",
     "Tarification en XAF, paiement Mobile Money, interface française. Conçu pour l'Afrique."),
    ("☁️", "Permanent",
     "Stockage décentralisé IPFS — le document existe à jamais, indépendamment de Certichain."),
]
for i, (icon, title, desc) in enumerate(points):
    y = Inches(1.75) + i * Inches(1.3)
    box(sl, Inches(0.55), y+Inches(0.05), Inches(0.5), Inches(0.5), TEAL)
    t(sl, icon, Inches(0.55), y+Inches(0.04), Inches(0.5), Inches(0.5),
      size=14, align=PP_ALIGN.CENTER, color=WHITE)
    t(sl, title, Inches(1.2), y, Inches(5.8), Inches(0.38),
      size=14, bold=True, color=NAVY)
    t(sl, desc,  Inches(1.2), y+Inches(0.38), Inches(5.8), Inches(0.75),
      size=11, color=DGRAY)

# Certificate mockup
bx2, by2, bw2, bh2 = Inches(7.5), Inches(1.65), Inches(5.25), Inches(5.35)
box(sl, bx2, by2, bw2, bh2, NAVY)
box(sl, bx2, by2, bw2, Pt(5), TEAL)
t(sl, "CERTICHAIN AFRICA", bx2, by2+Inches(0.2), bw2, Inches(0.38),
  size=9, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
t(sl, "DÉCERNE FIÈREMENT CE DIPLÔME À", bx2, by2+Inches(0.62), bw2, Inches(0.32),
  size=8, color=MGRAY, align=PP_ALIGN.CENTER)
t(sl, "Marie-Claire Nguema", bx2, by2+Inches(0.98), bw2, Inches(0.55),
  size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rule(sl, bx2+Inches(0.6), by2+Inches(1.65), bw2-Inches(1.2), TEAL, Pt(2))
t(sl, "Diplôme d'Ingénierie Informatique", bx2, by2+Inches(1.82), bw2, Inches(0.38),
  size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
t(sl, "Université de Yaoundé I  ·  Juin 2025", bx2, by2+Inches(2.28), bw2, Inches(0.32),
  size=9, color=MGRAY, align=PP_ALIGN.CENTER)
box(sl, bx2+bw2-Inches(1.0), by2+bh2-Inches(1.15), Inches(0.72), Inches(0.72), TEAL)
t(sl, "QR", bx2+bw2-Inches(1.0), by2+bh2-Inches(1.15), Inches(0.72), Inches(0.72),
  size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
t(sl, "Ancré sur Polygon · Hash SHA-256 · IPFS", bx2, by2+bh2-Inches(0.38), bw2, Inches(0.32),
  size=7.5, color=MGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, WHITE)
header(sl, "COMMENT ÇA MARCHE",
       "4 étapes — de l'émission à la vérification en 3 secondes")

steps = [
    ("1", "Institution\ns'inscrit",
     "Elle crée son compte,\nconfigure son profil\net choisit son plan."),
    ("2", "Émission du\ncertificat",
     "Elle saisit les données\nbénéficiaire. Certichain\ngénère un PDF + QR code."),
    ("3", "Ancrage\nblockchain",
     "L'empreinte SHA-256\nest ancrée sur Polygon\n(L2 Ethereum, ~$0,01)."),
    ("4", "Vérification\ninstantanée",
     "Le recruteur scanne\nle QR → résultat\nen 3 s, sans compte."),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.45) + i * Inches(3.2)
    oval(sl, x, Inches(1.8), Inches(1.1), Inches(1.1), NAVY)
    t(sl, num, x, Inches(1.8), Inches(1.1), Inches(1.1),
      size=30, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    if i < 3:
        t(sl, "→", x+Inches(1.15), Inches(2.1), Inches(2.1), Inches(0.5),
          size=22, color=TEAL, align=PP_ALIGN.CENTER)
    t(sl, title, x, Inches(3.18), Inches(3.0), Inches(0.72),
      size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    t(sl, desc,  x, Inches(3.95), Inches(2.95), Inches(2.1),
      size=10.5, color=DGRAY, align=PP_ALIGN.CENTER)

box(sl, 0, Inches(6.55), W, Inches(0.95), LGRAY)
t(sl, "STACK :  Polygon (L2 Ethereum)  ·  IPFS / Pinata  ·  SHA-256  ·  PDF + QR Code  ·  Flask REST API  ·  Mobile Money (MoMo)  ·  Python",
  Inches(0.6), Inches(6.65), Inches(12.2), Inches(0.55),
  size=10, color=DGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, WHITE)
header(sl, "OPPORTUNITÉ DE MARCHÉ",
       "Marché mondial de $6,4 Mds · 22 % TCAC · 54 pays africains")

# Concentric circles (TAM / SAM / SOM)
cx = Inches(3.3)
cy = Inches(4.6)
circle_data = [
    (Inches(4.4), Inches(4.4), RGBColor(0xE0, 0xEA, 0xFF),
     "$6,4 Mds", "TAM — Marché mondial\ncredentials numériques\n(2028, +22% TCAC)"),
    (Inches(3.0), Inches(3.0), RGBColor(0xCC, 0xF0, 0xEB),
     "$420 M*", "SAM — Institutions\nd'Afrique francophone"),
    (Inches(1.6), Inches(1.6), TEAL,
     "$18 M*", "SOM Y1–Y3"),
]
for (cw, ch, fill, val, lbl) in circle_data:
    ox = cx - cw/2; oy = cy - ch/2
    oval(sl, ox, oy, cw, ch, fill)

for (cw, ch, fill, val, lbl) in circle_data:
    ox = cx - cw/2; oy = cy - ch/2
    s = cw.inches
    tb = sl.shapes.add_textbox(ox+Pt(4), oy+Pt(4), cw-Pt(8), ch-Pt(8))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = val
    run.font.name = 'Calibri'; run.font.size = Pt(max(8.5, min(15, s*9)))
    run.font.bold = True; run.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = lbl.split('\n')[0]
    r2.font.name = 'Calibri'; r2.font.size = Pt(max(7, min(9.5, s*7)))
    r2.font.color.rgb = DGRAY

# Right stats
stats_r = [
    ("2 000+",    "universités en Afrique",             NAVY),
    ("50 000+",   "lycées & centres de formation",       NAVY),
    ("22 % TCAC", "marché credentials numériques 2023–28*", TEAL),
    ("12 M",      "diplômés africains / an*",             GOLD),
    ("15–30 %",   "taux de fraude documentaire estimé*",  RGBColor(0x9B, 0x1C, 0x1C)),
]
for i, (val, lbl, col) in enumerate(stats_r):
    y = Inches(1.72) + i * Inches(1.08)
    box(sl, Inches(7.15), y+Inches(0.05), Pt(5), Inches(0.65), col)
    t(sl, val, Inches(7.45), y,             Inches(5.3), Inches(0.48),
      size=21, bold=True, color=col)
    t(sl, lbl, Inches(7.45), y+Inches(0.45), Inches(5.3), Inches(0.45),
      size=10, color=DGRAY)

footnote(sl, "* Sources estimatives — UNESCO 2023, HolonIQ EdTech Africa, AfDB Education Report [PLACEHOLDER]")


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, WHITE)
header(sl, "MODÈLE DE REVENUS",
       "SaaS mensuel · Mobile Money XAF · 4 niveaux · −20 % annuel")

plans = [
    ("Découverte",  "Gratuit",      "",           "• 5 certs / mois\n• SHA-256 + QR code\n• Blockchain Polygon\n• Tableau de bord",
     BGRAY,                           DGRAY,  DGRAY),
    ("Starter",     "19 900 XAF",   "~$33 / mois","• 50 certs / mois\n• 3 types de certificat\n• Page vérif. publique\n• Support email 48h",
     RGBColor(0xEF, 0xF6, 0xFF),      NAVY,   NAVY),
    ("Académique",  "49 900 XAF",   "~$83 / mois","• 300 certs / mois\n• Email au bénéficiaire\n• Stockage IPFS permanent\n• Analytics · Support 24h",
     RGBColor(0xF0, 0xFD, 0xFA),      RGBColor(0x0D,0x7A,0x72), RGBColor(0x0D,0x7A,0x72)),
    ("Entreprise",  "Sur devis",    "volume",     "• Illimité + API REST\n• Multi-établissements\n• SSO / LMS · Webhooks\n• SLA 99,9 %",
     NAVY,                            GOLD,   GOLD),
]
accents = [MGRAY, NAVY, TEAL, GOLD]
for i, (name, price, sub, features, bg, title_c, price_c) in enumerate(plans):
    px = Inches(0.42) + i * Inches(3.22)
    pw = Inches(3.05)
    box(sl, px, Inches(1.75), pw, Inches(5.25), bg)
    box(sl, px, Inches(1.75), pw, Pt(5), accents[i])
    t(sl, name,     px+Inches(0.22), Inches(1.9),  pw-Inches(0.35), Inches(0.38),
      size=13, bold=True, color=title_c)
    t(sl, price,    px+Inches(0.22), Inches(2.35), pw-Inches(0.35), Inches(0.55),
      size=21, bold=True, color=price_c)
    if sub:
        t(sl, sub,  px+Inches(0.22), Inches(2.92), pw-Inches(0.35), Inches(0.28),
          size=9.5, color=MGRAY if bg != NAVY else RGBColor(0xCB,0xD5,0xE1))
    t(sl, features, px+Inches(0.22), Inches(3.28), pw-Inches(0.35), Inches(2.55),
      size=10.5, color=WHITE if bg == NAVY else DGRAY)

t(sl, "💡  Facturation annuelle : −20 % sur tous les plans payants  ·  Paiement par Mobile Money (Orange Money, MTN MoMo)",
  Inches(0.55), Inches(7.06), Inches(12.3), Inches(0.38),
  size=11, color=NAVY)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — GO-TO-MARKET
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, WHITE)
header(sl, "GO-TO-MARKET",
       "Afrique francophone d'abord · Expansion continentale · Leadership pan-africain")

phases = [
    ("Phase 1\n2025 – 2026", "Afrique francophone\nOuest & Centre",
     "• Cameroun, Côte d'Ivoire, Sénégal\n• Contact direct universités & grandes écoles\n• Partenariats ministères de l'Éducation\n• Programme démo gratuite 30 jours\n• Cible : 100 institutions payantes",
     NAVY,                           WHITE, TEAL),
    ("Phase 2\n2026 – 2027", "Expansion Afrique\ndu Nord & Anglophone",
     "• Maroc, Tunisie, Algérie\n• Ghana, Nigeria, Kenya\n• Programme revendeurs locaux\n• Intégration LMS régionaux\n• Cible : 500 institutions payantes",
     RGBColor(0x0D, 0x44, 0x60),     WHITE, TEAL),
    ("Phase 3\n2027+",       "Leadership\npan-africain",
     "• 54 pays africains\n• API gouvernementale vérification\n• Partenariats entreprises RH\n• Licences institutionnelles État\n• Cible : 2 000+ institutions",
     RGBColor(0x78, 0x35, 0x0F),     WHITE, GOLD),
]
for i, (phase, geo, bullets, bg, txt, acc) in enumerate(phases):
    px = Inches(0.45) + i * Inches(4.28)
    pw = Inches(4.0)
    box(sl, px, Inches(1.75), pw, Inches(5.1), bg)
    box(sl, px, Inches(1.75), pw, Pt(5), acc)
    t(sl, phase,   px+Inches(0.25), Inches(1.88), pw-Inches(0.4), Inches(0.72),
      size=14, bold=True, color=acc)
    t(sl, geo,     px+Inches(0.25), Inches(2.65), pw-Inches(0.4), Inches(0.55),
      size=12, bold=True, color=txt)
    t(sl, bullets, px+Inches(0.25), Inches(3.25), pw-Inches(0.4), Inches(2.65),
      size=10.5, color=txt)

box(sl, 0, Inches(6.85), W, Inches(0.65), LGRAY)
t(sl, "CANAUX :  Vente directe BtoI  ·  Partenariats institutionnels  ·  Marketing digital  ·  Démos gratuites  ·  Bouche-à-oreille inter-universités",
  Inches(0.6), Inches(6.92), Inches(12.2), Inches(0.42),
  size=10, color=DGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — COMPETITION
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, WHITE)
header(sl, "CONCURRENCE & DIFFÉRENCIATION",
       "La seule solution blockchain native pour l'Afrique francophone")

headers_row = ["Critère", "Papier / email", "Blockcerts\n(global)", "Accredible\n(global)", "Certichain Africa ✓"]
col_ws = [Inches(3.15), Inches(2.12), Inches(2.12), Inches(2.12), Inches(2.5)]
xs = [Inches(0.42)]
for w in col_ws[:-1]:
    xs.append(xs[-1] + w + Inches(0.04))

# Header row
hdr_cols = [NAVY, NAVY, NAVY, NAVY, TEAL]
for i, (hdr, cw, cx2) in enumerate(zip(headers_row, col_ws, xs)):
    box(sl, cx2, Inches(1.75), cw, Inches(0.55), hdr_cols[i])
    t(sl, hdr, cx2+Inches(0.1), Inches(1.78), cw-Inches(0.12), Inches(0.55),
      size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Paiement Mobile Money",         "✗", "✗",       "✗",       "✓"),
    ("Tarification en XAF",           "✗", "✗",       "✗",       "✓"),
    ("Interface français",            "✗", "Partiel",  "Non",     "✓"),
    ("Blockchain L2 bas coût",        "✗", "✓",        "✗",       "✓"),
    ("Stockage IPFS permanent",       "✗", "✓",        "✗",       "✓"),
    ("Vérification publique sans login","Partiel","✓", "✓",       "✓"),
    ("Support local Afrique",         "✗", "✗",        "✗",       "✓"),
]
row_bgs = [BGRAY, WHITE]
for r, (feat, *vals) in enumerate(rows):
    ry = Inches(2.34) + r * Inches(0.545)
    for i, (cx2, cw) in enumerate(zip(xs, col_ws)):
        rbg = RGBColor(0xF0,0xFD,0xFA) if i == 4 else row_bgs[r % 2]
        box(sl, cx2, ry, cw, Inches(0.52), rbg)
        txt = feat if i == 0 else vals[i-1]
        is_ok = (txt == "✓"); is_no = (txt == "✗")
        col = TEAL if is_ok else (RGBColor(0xCB,0xD5,0xE1) if is_no else DGRAY)
        al  = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
        t(sl, txt, cx2+Inches(0.1), ry+Inches(0.1), cw-Inches(0.12), Inches(0.35),
          size=11 if (i != 0 and (is_ok or is_no)) else 10,
          bold=(i == 0 or i == 4), color=col, align=al)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — TRACTION
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, WHITE)
header(sl, "TRACTION & JALONS",
       "Produit live · Paiements intégrés · Acquisition clients en cours")

milestones = [
    (True,  "MVP développé & déployé",           "Plateforme complète : émission, dashboard, tableau statistiques", "2024 Q4"),
    (True,  "4 types de certificats",            "Diplôme · Certification · Badge · Cyber/Tech Diplôme",          "2025 Q1"),
    (True,  "Ancrage blockchain Polygon",        "Hash SHA-256 on-chain + stockage IPFS via Pinata",              "2025 Q1"),
    (True,  "Intégration Mobile Money (MoMo)",  "Paiement XAF — cycles mensuel & annuel",                        "2025 Q2"),
    (True,  "Page vérification publique",        "Vérification par ID, hash ou fichier — sans compte requis",     "2025 Q2"),
    (False, "[X] institutions en bêta  [PLACEHOLDER]",  "Pilote avec établissements partenaires",                "2025 Q3"),
    (False, "[X] certificats émis  [PLACEHOLDER]",       "Volume en cours de validation product-market fit",     "2025 Q3"),
]
for i, (done, title, desc, date) in enumerate(milestones):
    y = Inches(1.85) + i * Inches(0.695)
    dc = TEAL if done else MGRAY
    oval(sl, Inches(0.48), y+Inches(0.06), Inches(0.4), Inches(0.4), dc)
    t(sl, "✓" if done else "⋯", Inches(0.48), y+Inches(0.04), Inches(0.4), Inches(0.4),
      size=10, bold=True, color=WHITE if done else DGRAY, align=PP_ALIGN.CENTER)
    t(sl, title, Inches(1.08), y,             Inches(8.0), Inches(0.32),
      size=12.5, bold=True, color=NAVY if done else MGRAY)
    t(sl, desc,  Inches(1.08), y+Inches(0.33), Inches(8.0), Inches(0.3),
      size=10, color=DGRAY)
    t(sl, date,  Inches(11.3), y, Inches(1.5), Inches(0.32),
      size=10, color=MGRAY, align=PP_ALIGN.RIGHT)

footnote(sl, "Remplacer les entrées [PLACEHOLDER] par vos chiffres réels avant toute présentation investisseur.")


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — TEAM
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, WHITE)
header(sl, "ÉQUIPE",
       "Des bâtisseurs engagés pour la transformation numérique africaine")

team = [
    ("[NOM FONDATEUR]",   "Fondateur / CEO",
     "[Ex. : Ingénieur logiciel · 8 ans en développement blockchain\nMBA Université de Yaoundé · Speaker EdTech Afrique]"),
    ("[NOM CO-FONDATEUR]","CTO / Co-fondateur",
     "[Ex. : Dev full-stack · Smart contracts Ethereum/Polygon\nAncien ingénieur [entreprise tech]  · Contributeur open-source]"),
    ("[NOM CMO]",         "CMO / Growth",
     "[Ex. : Marketing EdTech Afrique · Réseau 50+ universités\nAncien [entreprise] · Expert acquisition BtoI Afrique]"),
    ("[NOM CONSEILLER]",  "Conseiller stratégique",
     "[Ex. : Directeur UNESCO Africa Education\nRéseau gouvernemental · Expert politiques éducatives]"),
]
for i, (name, role, bio) in enumerate(team):
    tx = Inches(0.45) + (i % 2) * Inches(6.4)
    ty = Inches(1.82) + (i // 2) * Inches(2.55)
    oval(sl, tx, ty, Inches(1.0), Inches(1.0), LGRAY)
    t(sl, "👤", tx, ty, Inches(1.0), Inches(1.0), size=26, align=PP_ALIGN.CENTER, color=MGRAY)
    t(sl, name, tx+Inches(1.2), ty,             Inches(4.7), Inches(0.42),
      size=15, bold=True, color=NAVY)
    t(sl, role, tx+Inches(1.2), ty+Inches(0.42), Inches(4.7), Inches(0.38),
      size=12, bold=True, color=TEAL)
    t(sl, bio,  tx+Inches(1.2), ty+Inches(0.85), Inches(4.9), Inches(0.85),
      size=10, color=DGRAY)

t(sl, "⚠  Remplacer les [NOM] et [descriptions] par les vraies informations de l'équipe avant présentation",
  Inches(0.55), Inches(7.08), Inches(12.3), Inches(0.35),
  size=9, color=RGBColor(0xD9, 0x77, 0x06), italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — FINANCIAL PROJECTIONS
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, WHITE)
header(sl, "PROJECTIONS FINANCIÈRES",
       "Trajectoire vers $2,1 M ARR en 5 ans  [projections indicatives — PLACEHOLDER]", 22)

# Bar chart
chart_data = ChartData()
chart_data.categories = ['2025', '2026', '2027', '2028', '2029']
chart_data.add_series('ARR (USD)', (24_000, 120_000, 360_000, 936_000, 2_100_000))
chart_shape = sl.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.42), Inches(1.72), Inches(7.6), Inches(4.9),
    chart_data
)
ch = chart_shape.chart
ch.has_legend = False
ch.has_title  = False
try:
    ser = ch.series[0]
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = NAVY
except Exception:
    pass

# Table
table_rows = [
    ("Année", "Institutions", "ARR USD", "ARR (XAF ~)"),
    ("2025*", "50",    "$24 K",     "~14 M XAF"),
    ("2026*", "200",   "$120 K",    "~72 M XAF"),
    ("2027*", "500",   "$360 K",   "~216 M XAF"),
    ("2028*", "1 200", "$936 K",   "~562 M XAF"),
    ("2029*", "2 500", "$2,1 M", "~1,26 Mds XAF"),
]
tcol_ws = [Inches(1.12), Inches(1.22), Inches(1.22), Inches(1.45)]
tx0     = Inches(8.35)
for r, row in enumerate(table_rows):
    ty = Inches(1.72) + r * Inches(0.77)
    for c_i, (val, cw) in enumerate(zip(row, tcol_ws)):
        cx2 = tx0 + sum(tcol_ws[:c_i]) + Inches(0.04)*c_i
        rbg = NAVY if r == 0 else (BGRAY if r % 2 == 0 else WHITE)
        box(sl, cx2, ty, cw, Inches(0.74), rbg)
        t(sl, val, cx2+Inches(0.06), ty+Inches(0.18), cw-Inches(0.08), Inches(0.4),
          size=10, bold=(r == 0), color=WHITE if r == 0 else DGRAY,
          align=PP_ALIGN.CENTER)

footnote(sl, "* Projections indicatives — basées sur taux d'adoption estimé. Hypothèse : ARPU $40/mois → $70/mois. À valider avec données réelles. [PLACEHOLDER]")


# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — THE ASK
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, NAVY)
box(sl, 0, 0, W, Pt(5), TEAL)
box(sl, 0, H-Pt(5), W, Pt(5), TEAL)
box(sl, 0, 0, Inches(0.38), H, TEAL)
oval(sl, W-Inches(4.5), -Inches(1.8), Inches(5.8), Inches(5.8), RGBColor(0x1A,0x3A,0x72))
oval(sl, -Inches(2), H-Inches(2.8), Inches(4), Inches(4), RGBColor(0x14,0x30,0x68))

t(sl, "NOTRE DEMANDE", Inches(0.7), Inches(0.52), Inches(10), Inches(0.42),
  size=10, bold=True, color=TEAL)
t(sl, "Nous levons [MONTANT]  [PLACEHOLDER]",
  Inches(0.7), Inches(1.05), Inches(11.5), Inches(0.82),
  size=36, bold=True, color=WHITE)
t(sl, "Pour accélérer notre déploiement en Afrique francophone et valider le product-market fit à 200+ institutions",
  Inches(0.7), Inches(2.0), Inches(10.5), Inches(0.68),
  size=14, color=RGBColor(0xCB,0xD5,0xE1))

rule(sl, Inches(0.7), Inches(2.82), Inches(3.2), TEAL, Pt(4))

funds = [
    ("40 %", "Ingénierie & Infrastructure",
     "Développement API, scalabilité, sécurité, nouvelles intégrations LMS"),
    ("30 %", "Ventes & Marketing",
     "Équipe commerciale terrain, partenariats institutions, présence digitale Afrique"),
    ("20 %", "Opérations & Expansion",
     "Bureau local (Douala / Abidjan), équipe support, programme onboarding"),
    ("10 %", "Juridique & Conformité",
     "Enregistrement, RGPD Afrique, contrats institutionnels, IP"),
]
for i, (pct, title, desc) in enumerate(funds):
    y = Inches(3.08) + i * Inches(1.02)
    box(sl, Inches(0.7), y, Inches(1.15), Inches(0.68), TEAL)
    t(sl, pct, Inches(0.7), y, Inches(1.15), Inches(0.68),
      size=19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    t(sl, title, Inches(2.05), y+Inches(0.04), Inches(5.5), Inches(0.32),
      size=12.5, bold=True, color=WHITE)
    t(sl, desc,  Inches(2.05), y+Inches(0.37), Inches(9.5), Inches(0.3),
      size=10, color=MGRAY)

t(sl, "Runway cible : 18 mois  ·  Objectif : 200 institutions payantes  ·  ARR cible : $120 K",
  Inches(0.7), Inches(7.1), Inches(12.2), Inches(0.35),
  size=11, color=GOLD)


# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — VISION & CONTACT
# ══════════════════════════════════════════════════════════════════
sl = new_slide()
box(sl, 0, 0, W, H, NAVY)
box(sl, 0, 0, Inches(0.38), H, TEAL)
oval(sl, -Inches(1.5), -Inches(0.8), Inches(4.8), Inches(4.8), RGBColor(0x14,0x30,0x68))
oval(sl, W-Inches(3.5), H-Inches(2.2), Inches(4.2), Inches(4.2), RGBColor(0x78,0x35,0x0F))

t(sl, "NOTRE VISION", Inches(0.78), Inches(1.35), Inches(12), Inches(0.42),
  size=10, bold=True, color=TEAL)
t(sl, "Devenir l'infrastructure de confiance\ndes diplômes et certifications\npour toute l'Afrique.",
  Inches(0.78), Inches(1.88), Inches(11.2), Inches(2.55),
  size=33, bold=True, color=WHITE)

rule(sl, Inches(0.78), Inches(4.68), Inches(3.8), TEAL, Pt(4))

t(sl, '"Un diplôme africain — vérifié en 3 secondes, partout dans le monde."',
  Inches(0.78), Inches(4.88), Inches(10.5), Inches(0.65),
  size=15.5, color=RGBColor(0xCB,0xD5,0xE1), italic=True)

contacts = [
    ("🌐", "certichain.africa"),
    ("✉️", "contact@certichain.africa"),
    ("📍", "Douala, Cameroun"),
]
for i, (icon, txt) in enumerate(contacts):
    t(sl, f"{icon}  {txt}", Inches(0.78), Inches(5.72)+i*Inches(0.52),
      Inches(6.5), Inches(0.44), size=12.5, color=RGBColor(0xCB,0xD5,0xE1))

t(sl, "Merci",      Inches(8.4), Inches(2.85), Inches(4.5), Inches(1.15),
  size=58, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
t(sl, "Questions ?", Inches(8.4), Inches(4.1), Inches(4.5), Inches(0.6),
  size=19, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
out_dir  = r"C:\Users\AB SOLUTIOINS\Desktop"
out_pptx = os.path.join(out_dir, "Certichain_PitchDeck.pptx")
prs.save(out_pptx)
print(f"[OK] PPTX saved: {out_pptx}")

# Try PDF export via PowerPoint COM (requires MS Office)
try:
    import comtypes.client
    ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
    ppt_app.Visible = True
    deck = ppt_app.Presentations.Open(out_pptx, ReadOnly=True)
    out_pdf = os.path.join(out_dir, "Certichain_PitchDeck.pdf")
    deck.SaveAs(out_pdf, 32)   # 32 = ppSaveAsPDF
    deck.Close()
    ppt_app.Quit()
    print(f"[OK] PDF saved: {out_pdf}")
except Exception as e:
    print(f"[--] PDF skip: {e}")
    print(f"     Open PPTX in PowerPoint then File > Export > Create PDF/XPS")
